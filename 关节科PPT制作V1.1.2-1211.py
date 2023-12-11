import tkinter as tk
from tkinter import messagebox 
from tkinter import filedialog
from tkinter import ttk
import tkinter.font

import sys
import os 
from fnmatch import fnmatch

from pathlib import Path

import pandas as pd
from pandas import Series,DataFrame

import random

import pptx
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor

from PIL import Image, ImageTk

from win32com.client import Dispatch
import win32com.client

#import shutil
import time


###########更新版本记得改版本号和时间！
version='1.1.1'
versiontime='2023-10-21'


class Impor():                                                                      ##导入和信息转换部分
    def __init__(self,master) :
        def xlsxpath():
            global Folderpath
            Folderpath = filedialog.askdirectory()
            for root, dirs, files in os.walk(Folderpath): 
                for file in files: 
                    pa = root+'/'+file              #pa = os.path.join(root, file)
                    if fnmatch(pa, "*.xlsx") and not '$' in pa:
                        return(pa, messagebox.askquestion('确认操作','找到手术excel表：'+file+',是否继续？'))      #函数里return后的语句不再执行

        def imp():                                                                                                                ##import
            pa,flag=xlsxpath()   
            if not flag:
                return 0
            df = pd.read_excel(pa, usecols=[4, 5, 6, 7, 8, 10, 11, 13, 14, 18], dtype={"登记号": object,'床号': str},skiprows=[0,1],engine='openpyxl')
            df["row"]=""
            df["左右"]=""
            df1 = pd.DataFrame(columns = df.columns.tolist())
            df2 = pd.DataFrame(columns = df.columns.tolist())
            df3 = pd.DataFrame(columns = df.columns.tolist())
            df4 = pd.DataFrame(columns = df.columns.tolist())
            
            df['床号']=df['床号'].str.replace('床','')
            df['性别']=df['性别'].str.replace('男','male').replace('女','female')
            df['年龄']=df['年龄'].str.replace('岁','ys').replace('月','ms')
            df['左右']=df['部位'].str.split('').str[1].replace('左','L-').replace('右','R-').replace('双','B-')
            for i in range(0,len(df['术者'])):
                if '黄威' in ''.join(df.loc[i,'术者']):
                    df.loc[i,'术者']='黄威'
                else:
                    df.loc[i,'术者']=(df.loc[i,'术者']).split(',')[0]
                    #print(df.loc[i,'术者'])
            df['病人病区']=df['病人病区'].str.replace('骨科一病区','').str.replace('骨科三病区','17F')
            df['床号']=df['病人病区']+df['床号']
            #print(df)

            i = 0
            #dignosis error
            ed1 = 0     #李守民
            ed2 = 0     #史国光
            ed3 = 0     #朱晨
            ed4 = 0     #黄威
            #surgery error
            es1 = 0     
            es2 = 0
            es3 = 0
            es4 = 0
            #术者计数器
            i1 = -1     
            i2 = -1
            i3 = -1 
            i4 = -1

            for ind, row in df.iterrows():
                ##诊断
                if '股骨颈骨折' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='FNF'
                
                elif '股骨粗隆间骨折' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='Intertrochanteric fracture'
            
                elif '股骨头' in df.loc[i,'诊断'] and '坏死' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='ONFH'
                
                elif '假体障碍' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='PJD'
                
                elif '假体植入感染' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='PJI'

                elif '膝关节置换术后疼痛' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='Periprosthetic Joint Pain'                
                
                elif '髋关节病' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='ONFH'            #warning
                    df.loc[i,'部位']='hip'

                elif '髋关节结核' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='Coxotuberculosis' 
                    df.loc[i,'部位']='hip'

                elif '膝关节病' in df.loc[i,'诊断'] or '膝骨关节' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='knee OA'    #warning
                    df.loc[i,'部位']='knee'
                
                elif '僵硬' in df.loc[i,'诊断'] or '强硬' in df.loc[i,'诊断']:
                    if '肘' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='Elbow Stiffness'           
                    if '膝' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='Knee Stiffness'  

                elif '肘管综合征' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='Cubital Tunnel Syndrome'           
                
                elif '肘关节病' in df.loc[i,'诊断'] or '肘关节关节病' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='Elbow OA'   #warning
                
                elif '囊肿' in df.loc[i,'诊断']:
                    if '腘窝' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='Popliteal Cyst' 
                    elif '肘关节' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='Elbow joint Cyst'
                        df.loc[i,'部位']='Elbow'          #肘关节囊肿术名为上肢软组织切除术，提前确认
                        df.loc[i,'row'] = 98
                    elif '膝关节' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='Knee Cyst'
                        df.loc[i,'部位']='Knee'
                        df.loc[i,'row'] = 298
                
                elif '肿物' in df.loc[i,'诊断']:
                    if '骨肿物' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='Bone Mass' 
                    elif '关节肿物' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='Joint Mass'
                    
                elif '瘤' in df.loc[i,'诊断']:
                    if '骨样骨瘤' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='Osteoid Osteoma'                    
                    elif '骨软骨瘤' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='Osteochondroma'
                    elif '血管瘤' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='Angioma'
                    else:
                        df.loc[i,'诊断']='Tumor'
                
                elif '假体松动' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='prosthetic loosening'
                
                elif '类风湿关节炎' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='RA'
                
                elif '半月板疾患' in df.loc[i,'诊断'] or '半月板损伤' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='Meniscus injury'
                    df.loc[i,'部位']='knee'
                    df.loc[i,'row'] = 298 
                
                elif '髋骨关节炎' in df.loc[i,'诊断'] or '髋骨关节病' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='hip OA'

                elif '假体周围骨折' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='periprosthetic fracture'
                
                elif '先天性髋关节发育不良' in df.loc[i,'诊断'] or '先天性髋关节半脱位' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='DDH'
                
                elif '取除骨折内固定装置' in df.loc[i,'诊断'] or '取出内固定装置' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='internal fracture Removal'

                elif '骨折术后' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='After fracture operation'

                elif '骨折' in df.loc[i,'诊断']:
                    if '肘' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='elbow fracture'
                    elif '肱骨' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='humerus fracture'
                    elif '髌骨' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='patella fracture'
                    elif '股骨' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='femoral fracture'

                elif '畸形' in df.loc[i,'诊断']:
                    if '股骨' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='leg Deformity'

                elif '脱位' in df.loc[i,'诊断']:
                    if '髌骨' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='Patellar dislocation'
                        df.loc[i,'部位']='knee'               #row在术式定，如关节镜 
                    elif '髋关节' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='hip dislocation'

                elif '盘状半月板' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='discoid meniscus'
                    df.loc[i,'部位']='knee'

                elif '血友病性关节炎' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='hemophilic arthritis'   

                elif '纤维结构不良' in df.loc[i,'诊断']:
                    if '股骨' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='femur osteofibrous dysplasia' 
                        df.loc[i,'row'] = 298
                    else:df.loc[i,'诊断']='osteofibrous dysplasia' 

                elif '术后感染' in df.loc[i,'诊断']:
                    if '膝' in df.loc[i,'诊断']:
                        df.loc[i,'部位']='knee'
                    df.loc[i,'诊断']='postoperative infection'

                elif '绒毛结节状滑膜炎' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='Villonodular synovitis'

                elif '膝滑囊炎' in df.loc[i,'诊断']:
                    df.loc[i,'诊断']='Knee bursitis'
                    df.loc[i,'部位']='knee'

                elif '肘骨关节病' in df.loc[i,'诊断']:
                    if '松解' in df.loc[i,'手术名称']:
                        df.loc[i,'诊断']='Elbow stiffness'        #否则就不动，留着医生自行判断

                elif '髌骨关节病' in df.loc[i,'诊断']:
                    if '脱位复位' in df.loc[i,'手术名称']:
                        df.loc[i,'诊断']='Patellar dislocation'
                        df.loc[i,'部位']='knee'
                    
                elif '游离体' in df.loc[i,'诊断']:
                    if '膝关节' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='Knee loose body'
                        df.loc[i,'部位']='knee'
                    else:
                        df.loc[i,'诊断']='Loose body'
                
                elif '修整' in df.loc[i,'诊断']:
                    if '截肢残端' in df.loc[i,'诊断']:
                        df.loc[i,'诊断']='Amputation stump'
                        if '大腿' in df.loc[i,'部位']:
                            df.loc[i,'部位']='leg'
                    else:
                        df.loc[i,'诊断']='Loose body'

                elif '黄威' in df.loc[i,'术者']:                          #error
                    es4 += 1
                elif '李守民' in df.loc[i,'术者']:
                    es1 += 1
                elif '史国光' in df.loc[i,'术者']:
                    es2 += 1
                else: es3 += 1 
                
                ##手术
                if ('单髁' in df.loc[i,'手术名称'] or '固定平台' in df.loc[i,'手术名称']) and '膝关节' in df.loc[i,'手术名称'] :
                    if '固定平台' in df.loc[i,'手术名称']:
                        df.loc[i,'手术名称']='UKA(Fixation)' 
                        df.loc[i,'部位']='knee'
                        df.loc[i,'row'] = 299
                    else: 
                        df.loc[i,'手术名称']='UKA'                                            #UKA
                        df.loc[i,'部位']='knee'
                        df.loc[i,'row'] = 300
                
                elif '髋关节置换' in df.loc[i,'手术名称'] or '全髋置换术' in df.loc[i,'手术名称']:
                    df.loc[i,'手术名称']='THA'                                                #THA
                    df.loc[i,'部位']='hip'
                    df.loc[i,'row'] = 600
                
                elif '全膝关节置换术' in df.loc[i,'手术名称']:                               
                    df.loc[i,'手术名称']='TKA'                                                #TKA
                    df.loc[i,'部位']='knee'
                    df.loc[i,'row'] = 500
                
                elif '股骨头置换术' in df.loc[i,'手术名称'] or '半髋置换术' in df.loc[i,'手术名称']:
                    df.loc[i,'手术名称']='Hemiarthroplasty'                                   #Hemiarthroplasty
                    df.loc[i,'部位']='hip'
                    df.loc[i,'row'] = 598
            
                elif '胫骨上端高位截骨术' in df.loc[i,'手术名称'] or '胫骨截骨术' in df.loc[i,'手术名称']:
                    df.loc[i,'手术名称']='HTO'                                                #HTO
                    df.loc[i,'部位']='knee'
                    df.loc[i,'row'] = 200

                elif '髌股关节表面置换术' in df.loc[i,'手术名称']:
                    df.loc[i,'手术名称']='PFJ'                                                #PFJ
                    df.loc[i,'部位']='knee'
                    df.loc[i,'row'] = 199

                elif '股骨下端截骨术' in df.loc[i,'手术名称']:
                    df.loc[i,'手术名称']='DFO'                                                #其他截骨术
                    df.loc[i,'部位']='knee'
                    df.loc[i,'row'] = 201
                
                elif '关节镜' in df.loc[i,'手术名称']:
                    df.loc[i,'手术名称']='Arthroscopy'                                        #Arthroscopy
                    df.loc[i,'row'] = 1
                    if '半月板' in df.loc[i,'手术名称'] or '膝' in df.loc[i,'手术名称']:   
                        df.loc[i,'部位']='knee'
                
                elif '翻修术' in df.loc[i,'手术名称']:                                
                    if '髋' in df.loc[i,'手术名称']:
                        df.loc[i,'手术名称']='Revision'                                       #Revision+髋判断
                        df.loc[i,'部位']='hip'
                        df.loc[i,'row'] = 698
                    if '膝' in df.loc[i,'手术名称']:
                        df.loc[i,'手术名称']='Revision'
                        df.loc[i,'部位']='knee'
                        df.loc[i,'row'] = 550
                    else:
                        df.loc[i,'手术名称']='Revision'                                   
                        if '黄威' in df.loc[i,'术者']:                          #error
                            es4 += 1
                        elif '李守民' in df.loc[i,'术者']:
                            es1 += 1
                        elif '史国光' in df.loc[i,'术者']:
                            es2 += 1
                        else: es3 += 1 
                
                elif '松解术' in df.loc[i,'手术名称']:                                    
                    if '肘' in df.loc[i,'手术名称']:
                        df.loc[i,'手术名称']='Release'                                        #Release
                        df.loc[i,'部位']='elbow'
                        df.loc[i,'row'] = 198
                    elif '膝' in df.loc[i,'手术名称']:
                        df.loc[i,'手术名称']='Release' 
                        df.loc[i,'部位']='knee'
                        df.loc[i,'row'] = 298.25
                    else:
                        df.loc[i,'手术名称']='Release'                              
                        if '黄威' in df.loc[i,'术者']:                          #error
                            es4 += 1
                        elif '李守民' in df.loc[i,'术者']:
                            es1 += 1
                        elif '史国光' in df.loc[i,'术者']:
                            es2 += 1
                        else: es3 += 1

                elif '切除术' in df.loc[i,'手术名称'] or '截除' in df.loc[i,'手术名称']:                                                   
                    if '腘窝' in df.loc[i,'手术名称']:                                        #Resection+腘窝囊肿/瘤段切除/肘关节囊肿/。。。
                        df.loc[i,'部位']='Knee'
                        df.loc[i,'row'] = 298
                    elif  '髋' in df.loc[i,'部位'] or '坐骨' in df.loc[i,'部位'] or df.loc[i,'部位'] =='hip':               
                        df.loc[i,'部位'] ='hip'
                        df.loc[i,'row'] = 596
                    elif '腿' in df.loc[i,'部位'] or '股骨' in df.loc[i,'部位'] : 
                        df.loc[i,'部位'] ='leg'
                        df.loc[i,'row'] = 298
                    elif '手' in df.loc[i,'部位'] : 
                        df.loc[i,'部位'] ='hand'
                        df.loc[i,'row'] = 90
                    elif '膝' in df.loc[i,'部位'] or 'knee' in df.loc[i,'部位']: 
                        df.loc[i,'部位'] ='knee'
                        df.loc[i,'row'] = 298                                                            
                    elif '肘' in df.loc[i,'部位'] or df.loc[i,'部位'] == 'elbow': 
                        df.loc[i,'部位'] ='elbow'
                        df.loc[i,'row'] = 98
                    elif '肱' in df.loc[i,'部位'] or df.loc[i,'部位'] == 'arm': 
                        df.loc[i,'部位'] ='arm'
                        df.loc[i,'row'] = 98
                    elif '瘤段' in df.loc[i,'手术名称']:
                        df.loc[i,'部位']='Tumor'
                    elif '血管瘤' in df.loc[i,'手术名称']:
                        df.loc[i,'部位']='Angioma'
                    elif '黄威' in df.loc[i,'术者']:                          #error
                        es4 += 1
                    elif '李守民' in df.loc[i,'术者']:
                        es1 += 1
                    elif '史国光' in df.loc[i,'术者']:
                        es2 += 1
                    else: es3 += 1
                    df.loc[i,'手术名称']='Resection'
                
                elif '矫形术' in df.loc[i,'手术名称'] or '复位术' in df.loc[i,'手术名称']:
                    if '髌骨' in df.loc[i,'手术名称']:                                        #Reduction+髌骨脱位
                        df.loc[i,'部位'] ='knee'
                        df.loc[i,'row'] = 298.5  
                    elif '股骨' in df.loc[i,'手术名称']:
                        df.loc[i,'部位'] ='leg'
                        df.loc[i,'row'] = 597  
                    elif '髋关节' in df.loc[i,'手术名称']:
                        df.loc[i,'部位'] ='hip'
                        df.loc[i,'row'] = 599  
                    elif '黄威' in df.loc[i,'术者']:                          #error
                        es4 += 1
                    elif '李守民' in df.loc[i,'术者']:
                        es1 += 1
                    elif '史国光' in df.loc[i,'术者']:
                        es2 += 1
                    else: es3 += 1
                    df.loc[i,'手术名称']='Reduction'

                elif '去除术' in df.loc[i,'手术名称'] or '取出术' in df.loc[i,'手术名称']: 
                    if '桡骨' in df.loc[i,'手术名称']:                                        #Removal（内固定装置去除术）
                        df.loc[i,'部位'] ='elbow'
                        df.loc[i,'row'] = 199
                    elif '锁骨' in df.loc[i,'手术名称']: 
                        df.loc[i,'部位'] ='shoulder'
                        df.loc[i,'row'] = 102
                    elif '股骨' in df.loc[i,'手术名称']: 
                        df.loc[i,'部位'] ='leg'
                        df.loc[i,'row'] = 199.5 
                    elif '髌骨' in df.loc[i,'手术名称']: 
                        df.loc[i,'部位'] ='patella'
                        df.loc[i,'row'] = 560
                    
                    df.loc[i,'手术名称']='Removal'

                elif '内固定术' in df.loc[i,'手术名称']:
                    if '肘' in df.loc[i,'手术名称']:                                          #Fixation（内固定）
                        df.loc[i,'部位'] ='elbow'
                        df.loc[i,'row'] = 104
                    if '股骨' in df.loc[i,'手术名称']:
                        df.loc[i,'部位'] ='leg'
                        df.loc[i,'row'] = 199.2
                    df.loc[i,'手术名称']='Fixation'

                elif '肘关节置换' in df.loc[i,'手术名称']:                                         
                    df.loc[i,'部位'] ='elbow'                                                 #TEA（肘关节置换）
                    df.loc[i,'row'] = 190
                    df.loc[i,'手术名称']='TEA'

                elif '肩关节置换修复术' in df.loc[i,'手术名称']:                                         
                    df.loc[i,'部位'] ='shoulder'                                              #TEA（肩关节置换）
                    df.loc[i,'row'] = 103
                    df.loc[i,'手术名称']='TSA'

                elif '融合术' in df.loc[i,'手术名称']:
                    if '膝' in df.loc[i,'手术名称']:                                          #Fusion(融合术)
                        df.loc[i,'部位'] ='knee' 
                        df.loc[i,'row'] =555
                    if '髋' in df.loc[i,'手术名称']:
                        df.loc[i,'部位'] ='hip' 
                        df.loc[i,'row'] =699                  
                    df.loc[i,'手术名称']='Fusion'   
                
                elif '修整术' in df.loc[i,'手术名称']:
                    if 'leg' in df.loc[i,'部位']:                                          #Fusion
                        df.loc[i,'row'] =556
                    df.loc[i,'手术名称']='Trimming' 

                elif '清创术' in df.loc[i,'手术名称']:
                    if 'leg' in df.loc[i,'部位']:                                          #Debridement
                        df.loc[i,'row'] =556
                    df.loc[i,'手术名称']='Debridement' 
                
                elif '黄威' in df.loc[i,'术者']:                          #error
                    es4 += 1
                elif '李守民' in df.loc[i,'术者']:
                    es1 += 1
                elif '史国光' in df.loc[i,'术者']:
                    es2 += 1
                else: es3 += 1
                
                if df.loc[i,'row']=='':
                    df.loc[i,'row']=1

                #分装表格
                if '黄威' in df.loc[i,'术者']:
                    i4 += 1
                    df4.loc[i4] = df.loc[i]
                else:
                    if '朱晨' in df.loc[i,'术者']:
                        i3 += 1
                        df3.loc[i3] = df.loc[i]
                    else:
                        if '李守民' in df.loc[i,'术者']:
                            i1 += 1
                            df1.loc[i1] = df.loc[i]
                        else:
                            if '史国光' in df.loc[i,'术者']:
                                        i2 += 1
                                        df2.loc[i2] = df.loc[i]

                i += 1


            df1.sort_values('row',inplace=True)
            df2.sort_values('row',inplace=True)
            df3.sort_values('row',inplace=True)
            df4.sort_values('row',inplace=True)
            #print(df3)                                         #调试用,可注释
            self.page.destroy()
            mex(df1,df2,df3,df4,i1,i2,i3,i4)

        global ppti           #ppt计数器
        ppti = 1

        def mex(df1,df2,df3,df4,i1,i2,i3,i4):                      #make excel：包括制作表格，以及病人选择和信息修改，并将结果回传；。
            if ppti <=4:
                if ppti ==1:
                    master.destroy()
                if i1 > -1:
                    df = df1
                    i = i1
                    i1 = -1
                elif i2 > -1:
                    df = df2
                    i = i2
                    i2 = -1
                elif i4 > -1:
                    df = df4
                    i = i4
                    i4 = -1
                elif i3 > -1:
                    df = df3
                    i = i3
                    i3 = -1
                else: 
                    conc()
                    return 0
            else:
                conc()
                return 0
            self.root = tk.Tk()
            self.root.title('早会PPT制作 V0.0.1')
            self.root.geometry('730x750+'+str(int(leng/16))+'+'+str(int(wid/16)))
            #print(df)                                         #调试用,可注释
            #return 0   
            exl={
                '1':['床号',5],
                '2':['登记号',11],
                '3':['姓名',8],
                '4':['性别',8],
                '5':['年龄',5],
                '6':['诊断',24],
                '7':['手术名称',17],
                '8':['部位',9],
                '9':['术者',5]
            }
            for c in range(1,10):                           #列
                name = tk.Entry(self.root,width=exl[str(c)][1],bg='green',fg='white')
                name.insert(0,exl[str(c)][0])
                name.grid(row=1,column=c,padx=2,pady=2)

                for r in range(0,i+1):
                    globals()[df.iloc[i]['术者']+str(r)+str(c)] = tk.StringVar()
                    locals()['entry'+str(r)+str(c)] = tk.Entry(self.root,textvariable=globals()[df.iloc[i]['术者']+str(r)+str(c)],width=exl[str(c)][1])
                    locals()['entry'+str(r)+str(c)].insert(0,df.iloc[r][exl[str(c)][0]])
                    if c== 6 or c == 7 or c==8:
                        if '\u4e00' <= df.iloc[r][exl[str(c)][0]] <= '\u9fa5':      #是汉字
                            locals()['entry'+str(r)+str(c)].configure(fg='red')
                    locals()['entry'+str(r)+str(c)].grid(row=r+2,column=c,padx=2,pady=2)
           
            warn=tk.Label(self.root, text="********请检查所有患者的诊断、手术名称和部位是否正确（是否出现红色中文）********",fg='green')
            #warn.tag_add('red', 33.0, 34.0)
            #warn.tag_config('red', foreground='red') 
            warn.grid(row = i+4,column=1,columnspan=8, pady =20)
            tk.Button(self.root, text="下一步", command=lambda:jump(df,i,df1,df2,df3,df4,i1,i2,i3,i4,exl)).grid(row = i+5,column =9, pady =10)  #不加lambda在加载前jump就执行，command为none；lambda为匿名函数
            #image = Image.open("logo.jpg")
            #new_image=image.resize((100,100)) 
            #pyt = ImageTk.PhotoImage(new_image)
            #pic= tk.Label(self.root,image=pyt)
            #pic.img=pyt
            #pic.grid(row = i+5,column =0,columnspan=3)
            self.root.mainloop()

        def jump(df,i,df1,df2,df3,df4,i1,i2,i3,i4,exl):
            global ppti               #在函数里修改外界变量要声明全局
            self.root.destroy()
            df= df.copy()       ##深拷贝，否则报错A value is trying to be set on a copy of a slice from a DataFrame 
            for c in range(1,10):
                for r in range(0,i+1):
                    df.iloc[r][exl[str(c)][0]]=globals()[df.iloc[i]['术者']+str(r)+str(c)].get()
            #print(df)
            pptx(df)
            ppti += 1
            mex(df1,df2,df3,df4,i1,i2,i3,i4)

        global text 
        text='所有患者ppt已制作完成！\n\n\n\n日志如下：\n\n'
        
        def conc():     #conclusion
            self.root = tk.Tk()
            self.root.title('早会PPT制作 V0.0.1')
            self.root.geometry('500x700+'+str(int(leng/16))+'+'+str(int(wid/16)))
            txt=tk.Text(self.root,width=70, height=30)
            txt.insert('insert', text)
            Desired_font = tkinter.font.Font( family = "Microsoft YaHei UI", 
                                 size = 10, 
                                 weight = "bold")
            txt.configure(font=Desired_font)
            txt.grid()

        def color(t):
            color={
                'Release':[112,48,160],
                'TEA':[112,48,160],
                'TSA':[112,48,160],
                'Arthroscopy':[146,208,80],
                'Fixation':[38,38,38],
                'HTO':[0,176,240],
                'DFO':[0,176,240],
                'UKA':[0,176,80],
                'UKA(Fixation)':[0,176,80],
                'TKA':[255,0,0],
                'THA':[38,38,38],
                'Reduction':[38,38,38],
                'Removal':[70,10,244],
                'Revision':[47,85,151],
                'Resection':[112,48,160],
                'Hemiarthroplasty':[89,89,89],
                'Fusion':[20,176,188],
                'Trimming':[3,37,205]
            }
            if t in color:
                return color[t]
            else:
                return [112,48,160]
            
        def pptx(df):
            #print(df)
            global text
            prs = Presentation()
            prs.slide_height = 6858000                                                          #设置ppt的高度
            prs.slide_width = 12192000                                                          #设置ppt的宽度,指定16:9
            i = 0
            pe = 0                                                                              #pic error
            for ind, row in df.iterrows():
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                slide.notes_slide.notes_text_frame.text = 'Powered by ChenMo Automation Lab'

                ##part_1 
                left, top, width, height = Cm(0), Cm(0), Cm(10), Cm(1.7)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame

                para = tf.paragraphs[0]                                                         # 新增段落
                para.text = df.iloc[i]['手术名称']                                               # 向段落写入文字
                #print(para.text,',')
                para.alignment = PP_ALIGN.CENTER                                                # 居中 # para.line_spacing = 1.5    # 1.5 倍的行距

                font = para.font
                font.name = 'Times New Roman'                                                   # 字体类型    # font.bold = True    # 加粗
                font.size = Pt(36)                                                              # 大小
                font.color.rgb=RGBColor(255,255,255)
                tf.vertical_anchor = MSO_ANCHOR.TOP                                             # txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                fill = txBox.fill
                fill.solid()                                                                    # 纯色填充
                fill.fore_color.rgb = RGBColor(color(df.iloc[i]['手术名称'])[0], color(df.iloc[i]['手术名称'])[1], color(df.iloc[i]['手术名称'])[2])

                ##part_2
                if 'L' in df.iloc[i]['左右']:
                    ori = 'left '
                elif 'R' in df.iloc[i]['左右']:
                    ori = 'right '
                else: ori = 'both '
                t = random.randint(1, 3)
                tu = 'y' 
                if  'fracture' in df.iloc[i]['诊断'] or 'FNF' in df.iloc[i]['诊断']:
                    tu = 'd'
                if t != 1:
                    tu =tu + 's'

                left, top, width, height = Cm(0), Cm(2), Cm(13), Cm(10)
                txBox2 = slide.shapes.add_textbox(left, top, width, height)
                tf2 = txBox2.text_frame                                                         #para2 = tf2.add_paragraph() 
                para2=tf2.paragraphs[0]
                para2.text = str(df.iloc[i]['床号']) +' '+str(df.iloc[i]['姓名'])+' '+str(df.iloc[i]['性别'])+' '+str(df.iloc[i]['年龄'])+' '+str(df.iloc[i]['登记号'])+ '\n'
                para2.text = para2.text + 'pain in ' + ori + str(df.iloc[i]['部位']) + ' for ' + str(t) + tu +'\n\n' 
                para2.text = para2.text + 'Lab: normal \n\n'
                para2.text = para2.text + 'Dignosis: ' + str(df.iloc[i]['左右'])+ str(df.iloc[i]['诊断']) + '\n\n'
                para2.text = para2.text + 'Surgical plan: ' + str(df.iloc[i]['左右']) + str(df.iloc[i]['手术名称']) + '\n\n'
                para2.text = para2.text + 'Surgery date: 2023.' + str(os.path.basename(Folderpath))

                font2 = para2.font
                font2.name = 'Times New Roman'
                font2.size = Pt(28) 
                
                
                ##pic import
                #print(str(df.iloc[i]['床号']),'+',len(str(df.iloc[i]['床号'])),'\n')
                pp='x'
                if os.path.exists(Folderpath+'/'+str(df.iloc[i]['登记号'])):
                    pp = Folderpath+'/'+str(df.iloc[i]['登记号'])                       #pic path
                elif os.path.exists(Folderpath+'/'+str(df.iloc[i]['床号'])):
                    pp = Folderpath+'/'+str(df.iloc[i]['床号'])
                elif os.path.exists(Folderpath+'/'+str(df.iloc[i]['姓名'])):
                    pp = Folderpath+'/'+str(df.iloc[i]['姓名'])
                elif os.path.exists(Folderpath+'/0'+str(df.iloc[i]['床号'])):
                    pp = Folderpath+'/0'+str(df.iloc[i]['床号'])
                elif len(str(df.iloc[i]['床号'])) > 1:
                    if os.path.exists(Folderpath+'/'+str(df.iloc[i]['床号'])[1]):
                        pp = Folderpath+'/'+str(df.iloc[i]['床号'][1])
                    else: 
                        pe += 1
                        #print('未找到患者图片：',df.iloc[i]['姓名'])
                        text +='未找到患者图片：'+df.iloc[i]['姓名']+'\n'
                        i +=1
                        continue
                else: 
                    pe += 1
                    #print('未找到患者图片：',df.iloc[i]['姓名'])
                    text +='未找到患者图片：'+df.iloc[i]['姓名']+'\n'
                    i +=1
                    continue
                pi = 0
                pic= pd.DataFrame(columns =['path','ratio'])

                
                for files in os.walk(pp):
                    left, top, width,= Cm(14), Cm(1), Cm(7)                      
                    
                    for file in files[2]:                
                        if file.endswith(".jpg") or file.endswith(".png"):                      
                            img = Image.open(files[0]+'/'+file)
                            pic.loc[pi] = [str(files[0]+'/'+file),img.size[1]/img.size[0]]
                            
                            pi += 1
                
                                                    ######slide宽度    33.87    cm，高     19.05       cm######
                if pi > 0:
                        if pi == 1:                                                                                 #髋正位
                            if pic.iloc[0]['ratio'] < 1.2 or df.iloc[i]['部位'] == 'hip':
                                left, top, width= Cm(16), Cm(1), Cm(16.5)       
                                if pic.iloc[0]['ratio'] > 1.2: 
                                    left, top, width= Cm(16), Cm(1), Cm(16.5 / pic.iloc[0]['ratio'])             
                            else:
                                left, top, width= Cm(21), Cm(0.8), Cm(6.5)    
                            pict = slide.shapes.add_picture(pic.iloc[0]['path'], left, top, width)
                        
                        elif pi == 2:
                            pic=pic.sort_values('ratio', ascending=False)                                           #肘正侧位（两张超宽）、正侧位（正常或超宽）或髋正位+全长
                            if pic.iloc[0]['ratio'] < 0.7 and pic.iloc[1]['ratio'] < 0.7:           ##两张肘片需要旋转90度
                                im = Image.open(pic.iloc[0]['path'])
                                im_rotate = im.transpose(Image.Transpose.ROTATE_270)
                                im_rotate.save(pic.iloc[0]['path']+'转.jpg')
                                left, top, width = Cm(33.87 - 0.25*2 - 14*pic.iloc[0]['ratio'] - 14*pic.iloc[1]['ratio']), Cm(3.8), Cm(14*pic.iloc[0]['ratio'])
                                pict = slide.shapes.add_picture(pic.iloc[0]['path']+'转.jpg', left, top, width)
                                im = Image.open(pic.iloc[1]['path'])
                                im_rotate = im.transpose(Image.Transpose.ROTATE_270)
                                im_rotate.save(pic.iloc[1]['path']+'转.jpg')
                                left,width = Cm(33.87 - 0.25 - 14*pic.iloc[1]['ratio']), Cm(14*pic.iloc[1]['ratio'])
                                pict = slide.shapes.add_picture(pic.iloc[1]['path']+'转.jpg', left, top, width)

                            elif pic.iloc[0]['ratio'] < 1.2 and pic.iloc[1]['ratio'] < 1.2:
                                left, top, width = Cm(33.87 - 0.25*2 - 10 / pic.iloc[0]['ratio']- 10 / pic.iloc[1]['ratio']), Cm(5), Cm(10 / pic.iloc[0]['ratio'])
                                pict = slide.shapes.add_picture(pic.iloc[0]['path'], left, top, width)
                                left,width = Cm(33.87 - 0.25- 10 / pic.iloc[1]['ratio']), Cm(10 / pic.iloc[1]['ratio'])
                                pict = slide.shapes.add_picture(pic.iloc[1]['path'], left, top, width)

                            elif pic.iloc[1]['ratio'] < 0.8:
                                left, top, width = Cm(33.87 - 0.25*2 -17.5 - 14 / pic.iloc[0]['ratio']), Cm(3.8), Cm(14 / pic.iloc[0]['ratio'])
                                pict = slide.shapes.add_picture(pic.iloc[0]['path'], left, top, width)
                                left,width = Cm(33.87 - 0.25 -17.5), Cm(17.5)
                                pict = slide.shapes.add_picture(pic.iloc[1]['path'], left, top, width)   

                            else:
                                left, top, width = Cm(33.87 - 0.25*2 - 14 / pic.iloc[0]['ratio'] - 14 / pic.iloc[1]['ratio']), Cm(3.8), Cm(14 / pic.iloc[0]['ratio'])
                                pict = slide.shapes.add_picture(pic.iloc[0]['path'], left, top, width)
                                left,width = Cm(33.87 - 0.25 - 14 / pic.iloc[1]['ratio']), Cm(14 / pic.iloc[1]['ratio'])
                                pict = slide.shapes.add_picture(pic.iloc[1]['path'], left, top, width)   
                            
                        elif pi == 3:                                                                               #1个全长+正侧位
                            pic=pic.sort_values('ratio', ascending=False)

                            left, top, width = Cm(33.87 - 0.25*3 - 6*2 - 14.8 / pic.iloc[0]['ratio']), Cm(3.8), Cm(14.8 / pic.iloc[0]['ratio'])
                            pict = slide.shapes.add_picture(pic.iloc[0]['path'], left, top, width)
                            left,width = Cm(33.87 - 0.25*2 - 6*2), Cm(6)
                            pict = slide.shapes.add_picture(pic.iloc[1]['path'], left, top, width)
                            left,width = Cm(33.87 - 0.25 - 6), Cm(6)
                            pict = slide.shapes.add_picture(pic.iloc[2]['path'], left, top, width)

                        elif pi == 4:
                            pic=pic.sort_values('ratio', ascending=True)                                           #2个磁共振+正位和侧位

                            left,top,width = Cm(33.87 - 7*2 - 0.25*2), Cm(0.25),Cm(7)
                            pict = slide.shapes.add_picture(pic.iloc[0]['path'], left, top, width)                    
                            left,width = Cm(33.87 - 7 - 0.25), Cm(7)
                            pict = slide.shapes.add_picture(pic.iloc[1]['path'], left, top, width)                    
                            left,top,width = Cm(33.87 - 7*2 - 0.25*2),Cm(0.5 + 7 * pic.iloc[1]['ratio']), Cm(7)
                            pict = slide.shapes.add_picture(pic.iloc[2]['path'], left, top, width)                    
                            left,height = Cm(33.87 - 7 - 0.25), Cm(7)
                            pict = slide.shapes.add_picture(pic.iloc[3]['path'], left, top, width)

                        elif pi == 5:
                            pic=pic.sort_values('ratio', ascending=True)                                            #1个全长，2个磁共振+正侧位
                            
                            left,top,width = Cm(33.87 - 6*2 - 0.25*3 - 14.8 / pic.iloc[4]['ratio']), Cm(3.25),Cm(14.8 / pic.iloc[4]['ratio'])
                            pict = slide.shapes.add_picture(pic.iloc[4]['path'], left, top, width)   

                            left,top,width = Cm(33.87 - 6*2 - 0.25*2), Cm(0.25),Cm(6)
                            pict = slide.shapes.add_picture(pic.iloc[0]['path'], left, top, width)                    
                            left,width = Cm(33.87 - 6 - 0.25), Cm(6)
                            pict = slide.shapes.add_picture(pic.iloc[1]['path'], left, top, width)                    
                            left,top,width = Cm(33.87 - 6*2 - 0.25*2), Cm(0.5 + 6 * pic.iloc[1]['ratio']), Cm(6)
                            pict = slide.shapes.add_picture(pic.iloc[2]['path'], left, top, width)                    
                            left,height = Cm(33.87 - 6 - 0.25), Cm(6)
                            pict = slide.shapes.add_picture(pic.iloc[3]['path'], left, top, width)
                        
                        else:
                            pic=pic.sort_values('ratio', ascending=True)                    
                            for pii in range(0, pi) :
                                left,top,width = Cm(16 + 6.5*(pii%3)) , Cm(0.5 + 7 * int(pii/3)) , Cm(6)
                                pict = slide.shapes.add_picture(pic.iloc[pii]['path'], left, top, width)
                i +=1
                    
        
            if pe > 0: 
                #print('图片导入错误：',pe,'请查询手术是否已取消')
                text +='未找到 '+df.iloc[i-1]['术者']+' 的患者影像如上，合计：'+str(pe)+'名。请查询手术是否已取消'+'\n\n\n'
            else: 
                #print('所有患者图片导入成功')
                text +=df.iloc[i-1]['术者']+' 的所有患者影像导入成功！'+'\n\n\n'

            if str(df.iloc[i-1]['术者']) == '黄威':
                prs.save(Folderpath+'/'+str(os.path.basename(Folderpath))+' '+str(df.iloc[i-1]['术者'])+'组.pptx')
            else:
                prs.save(Folderpath+'/'+str(os.path.basename(Folderpath))+' '+str(df.iloc[i-1]['术者'])+'主任组.pptx')
           # if os.path.isfile(Folderpath+'/'+str(os.path.basename(Folderpath))+' '+str(df.iloc[i-1]['术者'])+'.pptx'):
            #    os.remove(Folderpath+'/'+str(os.path.basename(Folderpath))+' '+str(df.iloc[i-1]['术者'])+'.pptx')
            #shutil.move(str(os.path.basename(Folderpath))+' '+str(df.iloc[i-1]['术者'])+'.pptx',Folderpath)         #调试时可注释掉

        def resource_path(relative_path):
            if getattr(sys, 'frozen', False): #是否Bundle Resource
                base_path = sys._MEIPASS
            else:
                base_path = os.path.abspath(".")
            return os.path.join(base_path, relative_path)

#------------------------------------------------------------------------------术后操作部分------------------------------------------------
        def aspic(num,fapa):                #after surgery picture
            pic= pd.DataFrame(columns =['path','ratio'])
            pi=0
            for files in os.walk(fapa):
                if num in files[0] and '术后' in files[0]:            ##偷个懒，最好是根据日期来选择存放照片的日期，否则存在恰好住院号在照片编号和一个人在同一周做了两次手术，也无法排除该天没做而在第二天做的情况 
                    for file in files[2]:                
                        if file.endswith(".jpg") or file.endswith(".png"):                      
                            img = Image.open(files[0]+'/'+file)
                            pic.loc[pi] = [str(files[0]+'/'+file),img.size[1]/img.size[0]]
                            pi+=1
            return pic,pi

        def move_slide(presentation,old_index, new_index):
            xml_slides = presentation.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[old_index])
            xml_slides.insert(new_index, slides[old_index])

        def afsur():
            paths = filedialog.askopenfilenames(title='请选择需要做术后的PPT文件(可多选)',  filetypes=[ ('PPT', '.ppt .pptx')])
            epi=0       #error picture i 未找到的患者
            for path in paths:   
                #到单个ppt
                prs = Presentation(path)
                fapa =Path(r'%s'%(path)).parent   #引用字符串内容，获得上级文件夹
                ppi = len(prs.slides)
                
                for pppi in range(0,ppi):
                    slide=prs.slides[pppi]
                    num=slide.shapes[1].text_frame.paragraphs[0].runs[0].text.split(' ')[-1]
                    if len(num) < 6:
                        for index,run in enumerate(slide.shapes[1].text_frame.paragraphs[0].runs):
                            if 'pain' in run.text:
                                num =slide.shapes[1].text_frame.paragraphs[0].runs[index-1].text.split(' ')[-1]
                                break

                    pic,pi=aspic(num,fapa)

                    sli_layout = prs.slide_layouts[6]
                    sli = prs.slides.add_slide(sli_layout)
                    #sli.notes_slide.notes_text_frame.text = 'Powered by ChenMo Automation Lab'
                    
                    left, top, width, height = Cm(0.5), Cm(0.5), Cm(3), Cm(3)
                    txBox = sli.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame                                                          
                    para=tf.paragraphs[0]
                    para.text ='术后'
                    font = para.font
                    font.name = 'Times New Roman'
                    font.size = Pt(28) 

                    if pi > 0:
                        if pi == 1:                                                                                 #髋正位
                            left, top, width= Cm(33.87 - 0.25 -18.4 / pic.iloc[0]['ratio']), Cm(0.25), Cm(18.4 / pic.iloc[0]['ratio'])    
                            pict = sli.shapes.add_picture(pic.iloc[0]['path'], left, top, width)
                        
                        elif pi == 2:   #偷懒不判断了，一般2张图的要么是没做全长，要么是肘的
                            pic=pic.sort_values('ratio', ascending=False)                                           
                            left, top, width = Cm(3), Cm(0.25), Cm(18.4 / pic.iloc[0]['ratio'])
                            pict = sli.shapes.add_picture(pic.iloc[0]['path'], left, top, width)
                            left,width = Cm(3 + 0.25 + 18.4 / pic.iloc[0]['ratio']), Cm(18.4 / pic.iloc[1]['ratio'])
                            pict = sli.shapes.add_picture(pic.iloc[1]['path'], left, top, width)   
                            
                        elif pi == 3:                                                                               #1个全长+正侧位
                            pic=pic.sort_values('ratio', ascending=False)
                            left, top, width = Cm(3), Cm(0.25), Cm(18.4 / pic.iloc[0]['ratio'])
                            pict = sli.shapes.add_picture(pic.iloc[0]['path'], left, top, width)
                            left,width = Cm(33.87 - 0.25*2 - 18.4 / pic.iloc[1]['ratio']- 18.4 / pic.iloc[2]['ratio']), Cm(18.4 / pic.iloc[1]['ratio'])
                            pict = sli.shapes.add_picture(pic.iloc[1]['path'], left, top, width)
                            left,width = Cm(33.87 - 0.25 - 18.4 / pic.iloc[2]['ratio']), Cm(18.4 / pic.iloc[2]['ratio'])
                            pict = sli.shapes.add_picture(pic.iloc[2]['path'], left, top, width)
                        
                        else:
                            pic=pic.sort_values('ratio', ascending=True)                    
                            for pii in range(0, pi) :
                                left,top,width = Cm(3 + 6.5*(pii%3)) , Cm(0.5 + 7 * int(pii/3)) , Cm(6)
                                pict = sli.shapes.add_picture(pic.iloc[pii]['path'], left, top, width)
                    
                    else : epi+=1
                    
                for pppi in range(0,ppi):
                    move_slide(prs,pppi+ppi,pppi*2+1)
                
                prs.save(path.replace('.pptx','-术后.pptx'))
            
            mx1 = messagebox.showinfo(title='提示', message=('所有术后片已制作完成！共',epi,'人术后片未找到'))


    ##----------------------------------------------分组法做术后-------------------------------
        def pmerge(alldirs,sway):                                                                                   #Microsoft法合并
            ppt = Dispatch('PowerPoint.Application')
            ppt.Visible = 1 
            ppt.DisplayAlerts = 0  
            #拿到所有的ppt地址！
            #alldirs=[r'D:\0课件文件\#研究生\交班ppt\py\1.pptx',r'D:\0课件文件\#研究生\交班ppt\py\2.pptx'] 
            pptA = ppt.Presentations.Open(alldirs[0].replace('/','\\'))
            for dir in alldirs[1:len(alldirs):1]:
                pptB = ppt.Presentations.Open(dir.replace('/','\\'))
                numB=len(pptB.Slides)
                for i in range(1,numB+1):
                    time.sleep(0.01)       ####如果报错没在剪切板里找到东西，则加个等待减减速
                    pptB.Slides(i).Copy()
                    pptA.Slides.Paste()
                pptB.Close()
            time.sleep(0.01)
            #设置合并后ppt的路径！
            pptA.SaveAs(sway) 
            pptA.Close() 
            ppt.Quit()  

        def wmerge(alldirs,sway):                                                                                   #WPS法合并
            ppt = win32com.client.Dispatch('KWPP.Application')
            #ppt.Visible = 1 
            #ppt.DisplayAlerts = 0  
            #拿到所有的ppt地址！
            #alldirs=[r'D:\0课件文件\#研究生\交班ppt\py\1.pptx']#,r'D:\0课件文件\#研究生\交班ppt\py\2.pptx'] 
            print(alldirs[0])
            pptA = ppt.Presentations.Open(alldirs[0])
            for dir in alldirs[1:len(alldirs):1]:
                pptB = ppt.Presentations.Open(dir)
                numB=len(pptB.Slides)
                for i in range(1,numB+1):
                    pptB.Slides(i).Copy()
                    time.sleep(0.01)
                    pptA.Slides.Paste()
                pptB.Close()
            #设置合并后ppt的路径！
            pptA.SaveAs(sway) 
            pptA.Close() 
            ppt.Quit()  
        
        def judgege(bed):           #judge group by excel
            df=pd.read_excel('床位分配（如果分组变动请修改）.xlsx',usecols=[0, 1])
            df['床号'] = df['床号'].astype(str)
            row=df.loc[df['床号'] == bed]   #搜索并提取
            try:
                name= str(row.iloc[:,1].iloc[0])
                return(name)
            except:
                try:
                    row=df.loc[df['床号'] == bed[1]]
                    name= str(row.iloc[:,1].iloc[0])
                    return(name)
                except:
                    return('x')       
     
        def mergeafpic(temppath,fapa):      #插入术后片
            epi = 0 
            prs = Presentation(temppath)
            ppi = len(prs.slides)
            for pppi in range(0,ppi):
                slide=prs.slides[pppi]
                try:
                    num=slide.shapes[1].text_frame.paragraphs[0].runs[0].text.split(' ')[-1]            #get num   ： 登记号
                except:
                    try:
                        for index,run in enumerate(slide.shapes[1].text_frame.paragraphs[0].runs):
                            if 'pain' in run.text:
                                num =slide.shapes[1].text_frame.paragraphs[0].runs[index-1].text.split(' ')[-1]
                                break
                    except:
                        for shape in slide.shapes:
                            try:
                                for index,run in enumerate(shape.text_frame.paragraphs[0].runs):
                                    if 'pain' in run.text:
                                        num =slide.shapes[1].text_frame.paragraphs[0].runs[index-1].text.split(' ')[-1]
                                        break
                            except:
                                continue

                if len(num) < 6:
                    for index,run in enumerate(slide.shapes[1].text_frame.paragraphs[0].runs):
                        if 'pain' in run.text:
                            num =slide.shapes[1].text_frame.paragraphs[0].runs[index-1].text.split(' ')[-1]
                            break

                pic,pi=aspic(num,fapa)

                sli_layout = prs.slide_layouts[6]
                sli = prs.slides.add_slide(sli_layout)
                #sli.notes_slide.notes_text_frame.text = 'Powered by ChenMo Automation Lab'
                
                left, top, width, height = Cm(0.5), Cm(0.5), Cm(3), Cm(3)
                txBox = sli.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame                                                          
                para=tf.paragraphs[0]
                para.text ='术后'
                font = para.font
                font.name = 'Times New Roman'
                font.size = Pt(28) 

                if pi > 0:
                    if pi == 1:                                                                                 #髋正位
                        left, top, width= Cm(33.87 - 0.25 -18.4 / pic.iloc[0]['ratio']), Cm(0.25), Cm(18.4 / pic.iloc[0]['ratio'])    
                        pict = sli.shapes.add_picture(pic.iloc[0]['path'], left, top, width)
                    
                    elif pi == 2:   #偷懒不判断了，一般2张图的要么是没做全长，要么是肘的
                        pic=pic.sort_values('ratio', ascending=False)                                           
                        left, top, width = Cm(3), Cm(0.25), Cm(18.4 / pic.iloc[0]['ratio'])
                        pict = sli.shapes.add_picture(pic.iloc[0]['path'], left, top, width)
                        left,width = Cm(3 + 0.25 + 18.4 / pic.iloc[0]['ratio']), Cm(18.4 / pic.iloc[1]['ratio'])
                        pict = sli.shapes.add_picture(pic.iloc[1]['path'], left, top, width)   
                        
                    elif pi == 3:                                                                               #1个全长+正侧位
                        pic=pic.sort_values('ratio', ascending=False)
                        left, top, width = Cm(3), Cm(0.25), Cm(18.4 / pic.iloc[0]['ratio'])
                        pict = sli.shapes.add_picture(pic.iloc[0]['path'], left, top, width)
                        left,width = Cm(33.87 - 0.25*2 - 18.4 / pic.iloc[1]['ratio']- 18.4 / pic.iloc[2]['ratio']), Cm(18.4 / pic.iloc[1]['ratio'])
                        pict = sli.shapes.add_picture(pic.iloc[1]['path'], left, top, width)
                        left,width = Cm(33.87 - 0.25 - 18.4 / pic.iloc[2]['ratio']), Cm(18.4 / pic.iloc[2]['ratio'])
                        pict = sli.shapes.add_picture(pic.iloc[2]['path'], left, top, width)
                    
                    else:
                        pic=pic.sort_values('ratio', ascending=True)                    
                        for pii in range(0, pi) :
                            left,top,width = Cm(3 + 6.5*(pii%3)) , Cm(0.5 + 7 * int(pii/3)) , Cm(6)
                            pict = sli.shapes.add_picture(pic.iloc[pii]['path'], left, top, width)
                
                else : epi+=1
                
            for pppi in range(0,ppi):
                move_slide(prs,pppi+ppi,pppi*2+1)

            prs.save(temppath)
            return epi            

        def afsurg():   #after surgery group分组法做术后

            paths = filedialog.askopenfilenames(title='请选择需要做术后的PPT文件(可多选)',  filetypes=[ ('PPT', '.ppt .pptx')])
            epi = 0   #error picture i 未找到的患者
            zcpaths=tuple(path for path in paths if '朱晨' in path)                  #拆分出朱晨path
            hwpaths=tuple(path for path in paths if '黄威' in path)                      #拆分黄威path
            progress=tkinter.Tk()
            #progress['height']=200
            #progress['width']=400
            pleng= progress.winfo_screenheight()
            pwid= progress.winfo_screenwidth()
            progress.geometry("300x100+"+str(int(pleng/2))+'+'+str(int(pwid/2)))
            progress.mainloop
            label=tkinter.Label(progress,text='制作进度:',font=('黑体', 10))
            label.place(x=10,y=60)
            #进度条
            progressbar=ttk.Progressbar(progress)#.grid(row = 8,column =1, pady =10)
            progressbar.place(x=80,y=60)
           
            progressbar['maximum']=6        #设置进度条最大值为6
            progressbar['length']=200       #设置进度条长度

            #合成PPT
            fapa =Path(r'%s'%(paths[0])).parent
            temppath=os.path.join(fapa,'#所有片子（本文件建议别删）.pptx')
            hwtemppath=os.path.join(fapa,'黄威组.pptx')
            zctemppath=os.path.join(fapa,'zctempt.pptx')
            ways=tk.messagebox.askquestion(title="选择辅助软件", message="是否使用WPS帮助合成PPT？（选否使用微软PPT帮助合成）\n\n\n※※注意！！WPS有一定概率合成失败，笔记本用户可以考虑插上电源后使用或者使用微软PPT合成")
            if ways == 'yes':
                #wmerge(paths,temppath)     #总表
                if not len(zcpaths) == 0:                       #有朱晨组
                    wmerge(zcpaths,zctemppath)
                    if not len(hwpaths) ==0:                    #有黄威组
                        wmerge(hwpaths,hwtemppath)
                        epi=mergeafpic(zctemppath,fapa)+mergeafpic(hwtemppath,fapa)
                        wmerge([zctemppath,hwtemppath],temppath)
                    else:
                        epi=mergeafpic(zctemppath,fapa)
                        wmerge(zctemppath,temppath)
                else:                                           #没有朱晨组
                    if not len(hwpaths) ==0:                    #有黄威组
                        wmerge(hwpaths,hwtemppath)
                        epi=mergeafpic(hwtemppath,fapa)
                        wmerge(hwtemppath,temppath)
                    else:
                        mx1 = messagebox.showinfo(title='提示', message=('没有可以制作的术前片！'))                  
            else:
                #pmerge(paths,temppath)     #总表
                if not len(zcpaths) == 0:                       #有朱晨组
                    pmerge(zcpaths,zctemppath)
                    if not len(hwpaths) ==0:                    #有黄威组
                        pmerge(hwpaths,hwtemppath)
                        epi=mergeafpic(zctemppath,fapa)+mergeafpic(hwtemppath,fapa)
                        pmerge([zctemppath,hwtemppath],temppath)
                    else:
                        epi=mergeafpic(zctemppath,fapa)
                        pmerge(zctemppath,temppath)
                else:                                           #没有朱晨组
                    if not len(hwpaths) ==0:                    #有黄威组
                        pmerge(hwpaths,hwtemppath)
                        epi=mergeafpic(hwtemppath,fapa)
                        pmerge(hwtemppath,temppath)
                    else:
                        mx1 = messagebox.showinfo(title='提示', message=('没有可以制作的术前片！')) 
                        #return 0 
            progressbar['value']=1
            #self.page.update()
            
            ##根据床号筛选
            for i in range(1,5):   
                prs=Presentation('C:/Users/cm/Documents/WeChat Files/wxid_33gkbmgoatrb22/FileStorage/File/2023-10/2023.10.23-2023.10.29术后/zctempt2.pptx')        #只有朱晨组的要筛选，黄威组的直接从前面导出来了
                ppi = len(prs.slides)
                jump=0
                ##对单个ppt
                for pppi in range(ppi//2):
                    pnum=(pppi-jump)*2
                    slide=prs.slides[pnum]
                    #print('t:',slide.shapes[0].text_frame.paragraphs[0].runs[0].text)
                    #到单个页面             
                    try:
                        bed=slide.shapes[1].text_frame.text.split(' ')[0]
                    except:
                        try:
                            t=slide.shapes[1].text_frame.paragraphs[0].runs[0].text     #其实不知道这后面的except存在还有没有意义
                            #print('1')
                        except:
                            t=slide.shapes[0].text_frame.paragraphs[0].runs[0].text
                            #print('2')
                        bed=t.split(' ')[0]
                    #print(bed)

                    
                    if not bed[-1].isdigit():
                        pi = 0  
                        bed='空'
                    
                    print(bed,judgege(bed))
                    if not str(i) in judgege(bed).replace('朱晨','1').replace('张贤祚','2').replace('李乾明','3').replace('禹德万','4'):    
                        rId = prs.slides._sldIdLst[pnum].rId
                        prs.part.drop_rel(rId)
                        del prs.slides._sldIdLst[pnum]
                        rId = prs.slides._sldIdLst[pnum].rId
                        prs.part.drop_rel(rId)
                        del prs.slides._sldIdLst[pnum]
                        jump+=1
                    

                if i==1:
                    prs.save(os.path.join(fapa,'朱晨组.pptx'))
                elif i==2:
                    prs.save(os.path.join(fapa,'张贤祚组.pptx'))
                elif i==3:
                    prs.save(os.path.join(fapa,'李乾明组.pptx'))
                elif i==4:
                    prs.save(os.path.join(fapa,'禹德万组.pptx'))
                #elif i==5:
                    #prs.save(os.path.join(fapa,'黄威组.pptx'))
                else:
                    print('zzzzzzzzzz')
                
                #print(jump)
                progressbar['value']=2+i
                progress.update()
            progress.destroy()
            #progress.mainloop
            mx1 = messagebox.showinfo(title='提示', message=('所有术后片已制作完成！共',epi,'人术后片未找到'))                
            os.remove(zctemppath)

 



################################################-----------------------------------------------------------------------
        #初始化界面
        filename = resource_path(os.path.join("icon","logo.jpg"))
        image = Image.open(filename)
        new_image=image.resize((100,100)) 
        pyt = ImageTk.PhotoImage(new_image)
        master.title("早会PPT制作 V"+version)
        
        global leng,wid
        leng= master.winfo_screenheight()
        wid= master.winfo_screenwidth()
        master.geometry("700x520+"+str(int(leng/16))+'+'+str(int(wid/16)))

        self.page= tk.Frame(master)
        self.page.pack()

        tk.Label(self.page,text='请点下方按钮选择存放手术单和病人影像的文件夹以开始：').grid(row = 1, column = 1,pady =20)
        tk.Button(self.page, text="选择要存放术前片的文件夹", command=imp).grid(row = 2,column =1, pady =10)
        tk.Label(self.page,text='使用说明：请将HIS系统导出的手术表（et格式）在WPS中转为xlsx格式，并将xlsx与存\n有单独患者影像的文件夹一起放在同个文件夹内，然后点按上方按钮后选择该文件夹。',fg = 'green').grid(row = 3,column =1, pady =10)
        tk.Label(self.page,text='注意事项：开始前请务必关闭excel表格！\n导片时请关注是否为骨科三病区患者，若是则需核对该患者是否真的在17F※※\n存放片子的文件夹名可以取“17F+床号”或登记号\n',fg='red').grid(row = 4,column =1)
        tk.Label(self.page,text='【因为程序是根据是否为三病区确定\n是否为17F，术前不整理好的话会影响\n术后分组】',fg='red').grid(row = 4,column =0)
        tk.Label(self.page,text='----------------------------------------').grid(row = 5,column =1)
        tk.Label(self.page,text='点击下方按钮选择要做术后ppt的方式（可多选）').grid(row = 6,column =1)
        tk.Button(self.page, text="按日期分装ppt", command=afsur).grid(row = 7,column =0, pady =10)
        tk.Button(self.page, text="按分组分装ppt（需要WPS或者微软PPT软件）", command=afsurg).grid(row = 7,column =1, pady =10)
        tk.Label(self.page,text='版本信息：Powered by ChenMo. Version: '+version+'，'+versiontime+'\n有任何使用问题和建议可联系作者：chenmomo@mail.ustc.edu.cn',font=('Arial',8)).grid(row = 9,column =1, pady =50)

        tk.Label(self.page,image=pyt).grid(row = 1,column =0)
        root.mainloop()



##----------------------------------main-----------------------------------
def resource_path(relative_path):
    if getattr(sys, 'frozen', False): #是否Bundle Resource
        base_path = sys._MEIPASS
    else:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)

if __name__ == '__main__':
    root = tk.Tk() 
    filename = resource_path(os.path.join("icon","ico.png"))
    root.iconphoto(True, tk.PhotoImage(file=filename))
    Impor(master = root)